# ppt_text_extractor.py

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches
from pptx.dml.color import RGBColor

def extract_text_from_ppt(ppt_path):
    prs = Presentation(ppt_path)
    extracted_data = []

    for slide_idx, slide in enumerate(prs.slides):
        slide_text_data = {"slide_index": slide_idx, "texts": []}
        for shape_idx, shape in enumerate(slide.shapes):
            # Handle text frames directly attached to shapes
            if shape.has_text_frame:
                text_frame = shape.text_frame
                if text_frame.text.strip(): # Check if text frame has non-empty text
                    for paragraph_idx, paragraph in enumerate(text_frame.paragraphs):
                        full_text = ""
                        for run_idx, run in enumerate(paragraph.runs):
                            full_text += run.text

                        if full_text.strip(): # Only add non-empty text
                            # Get font color safely
                            font_color_rgb = None
                            if paragraph.font.color.type is not None and paragraph.font.color.rgb is not None:
                                font_color_rgb = paragraph.font.color.rgb.rgb

                            # Get position and size safely
                            left, top, width, height = None, None, None, None
                            try:
                                left = shape.left.emu
                                top = shape.top.emu
                                width = shape.width.emu
                                height = shape.height.emu
                            except AttributeError: # Some shapes (like placeholders) might not have these directly
                                pass

                            slide_text_data["texts"].append({
                                "shape_index": shape_idx,
                                "paragraph_index": paragraph_idx,
                                "original_text": full_text,
                                "shape_type": shape.shape_type.name, # e.g., 'TEXT_BOX', 'AUTO_SHAPE'
                                "left": left,
                                "top": top,
                                "width": width,
                                "height": height,
                                "font_name": paragraph.font.name,
                                "font_size": paragraph.font.size.emu if paragraph.font.size else None,
                                "font_bold": paragraph.font.bold,
                                "font_italic": paragraph.font.italic,
                                "font_color_rgb": font_color_rgb,
                                "alignment": paragraph.alignment.name if paragraph.alignment else None
                            })
            # Handle tables
            elif shape.has_table:
                table = shape.table
                for row_idx, row in enumerate(table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        if cell.text.strip():
                            # Get font color safely for table cells
                            font_color_rgb = None
                            if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].font.color.type is not None and cell.text_frame.paragraphs[0].font.color.rgb is not None:
                                font_color_rgb = cell.text_frame.paragraphs[0].font.color.rgb.rgb

                            slide_text_data["texts"].append({
                                "shape_index": shape_idx,
                                "row_index": row_idx,
                                "col_index": col_idx,
                                "original_text": cell.text,
                                "shape_type": "TABLE_CELL",
                                "font_name": cell.text_frame.paragraphs[0].font.name if cell.text_frame.paragraphs else None,
                                "font_size": cell.text_frame.paragraphs[0].font.size.emu if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].font.size else None,
                                "font_bold": cell.text_frame.paragraphs[0].font.bold if cell.text_frame.paragraphs else None,
                                "font_italic": cell.text_frame.paragraphs[0].font.italic if cell.text_frame.paragraphs else None,
                                "font_color_rgb": font_color_rgb,
                                "alignment": cell.text_frame.paragraphs[0].alignment.name if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].alignment else None
                            })
        if slide_text_data["texts"]:
            extracted_data.append(slide_text_data)

    return extracted_data

if __name__ == "__main__":
    # Example usage (requires a dummy.pptx file for testing)
    # You would typically call this function from your main application logic
    # For demonstration, let's create a dummy PPTX file first

    prs = Presentation()
    slide_layout = prs.slide_layouts[1] # Title and Content layout

    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "Hello World!"
    content.text = "This is a test paragraph.\nAnother line of text.\n\nAnd a third one."

    # Add a table
    rows = 2
    cols = 2
    left = top = width = height = Inches(1.0)
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    cell = table.cell(0, 0)
    cell.text = "Header 1"
    cell = table.cell(0, 1)
    cell.text = "Header 2"
    cell = table.cell(1, 0)
    cell.text = "Data A"
    cell = table.cell(1, 1)
    cell.text = "Data B"

    # Add another slide with just a text box
    slide_layout_blank = prs.slide_layouts[6] # Blank layout
    slide2 = prs.slides.add_slide(slide_layout_blank)
    left = top = width = height = Inches(1.0)
    txBox = slide2.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "This is a standalone text box."

    dummy_ppt_path = "dummy.pptx"
    prs.save(dummy_ppt_path)
    print(f"Created dummy PPTX at {dummy_ppt_path}")

    # Now, extract text from the dummy PPTX
    extracted_texts = extract_text_from_ppt(dummy_ppt_path)
    import json
    print(json.dumps(extracted_texts, indent=4, ensure_ascii=False))

    # Clean up the dummy file
    import os
    os.remove(dummy_ppt_path)
    print(f"Removed dummy PPTX at {dummy_ppt_path}")


