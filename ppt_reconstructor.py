# ppt_reconstructor.py

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE_TYPE

def reconstruct_ppt(original_ppt_path, translated_data, output_ppt_path):
    prs = Presentation(original_ppt_path)

    for slide_data in translated_data:
        slide_idx = slide_data["slide_index"]
        slide = prs.slides[slide_idx]

        for text_info in slide_data["texts"]:
            shape_idx = text_info["shape_index"]
            translated_text = text_info["translated_text"]

            # Handle text in shapes (text boxes, placeholders)
            if text_info["shape_type"] != "TABLE_CELL":
                shape = slide.shapes[shape_idx]
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    paragraph_idx = text_info["paragraph_index"]

                    # Clear existing runs in the paragraph and add the translated text
                    if paragraph_idx < len(text_frame.paragraphs):
                        paragraph = text_frame.paragraphs[paragraph_idx]
                        # Clear existing runs
                        for i in range(len(paragraph.runs) - 1, -1, -1):
                            p = paragraph._p
                            p.remove(paragraph.runs[i]._r)

                        run = paragraph.add_run()
                        run.text = translated_text

                        # Apply original formatting if available
                        if text_info["font_name"]:
                            run.font.name = text_info["font_name"]
                        if text_info["font_size"]:
                            run.font.size = Pt(text_info["font_size"] / 914400 * 72) # Convert EMU to Pt
                        if text_info["font_bold"] is not None:
                            run.font.bold = text_info["font_bold"]
                        if text_info["font_italic"] is not None:
                            run.font.italic = text_info["font_italic"]
                        if text_info["font_color_rgb"] is not None:
                            # RGBColor expects r, g, b as separate arguments
                            # Assuming text_info["font_color_rgb"] is an integer like 0xRRGGBB
                            rgb_int = text_info["font_color_rgb"]
                            r = (rgb_int >> 16) & 0xFF
                            g = (rgb_int >> 8) & 0xFF
                            b = rgb_int & 0xFF
                            run.font.color.rgb = RGBColor(r, g, b)
                        if text_info["alignment"]:
                            # Convert string alignment to MSO_ALIGNMENT enum
                            # This requires a mapping from string to enum, or direct enum values
                            # For simplicity, we'll skip complex alignment for now or assume basic ones
                            pass # TODO: Implement alignment mapping

                        # Adjust shape size if text overflows (basic attempt)
                        # This is a complex problem, python-pptx doesn't have auto-fit like PowerPoint
                        # For now, we rely on PowerPoint's default auto-fit or manual adjustment
                        # text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT # This might not work as expected for all shapes

            # Handle text in table cells
            elif text_info["shape_type"] == "TABLE_CELL":
                shape = slide.shapes[shape_idx]
                if shape.has_table:
                    table = shape.table
                    row_idx = text_info["row_index"]
                    col_idx = text_info["col_index"]

                    if row_idx < len(table.rows) and col_idx < len(table.columns):
                        cell = table.cell(row_idx, col_idx)
                        cell.text = translated_text

                        # Apply formatting to the first paragraph in the cell
                        if cell.text_frame.paragraphs:
                            paragraph = cell.text_frame.paragraphs[0]
                            if text_info["font_name"]:
                                paragraph.font.name = text_info["font_name"]
                            if text_info["font_size"]:
                                paragraph.font.size = Pt(text_info["font_size"] / 914400 * 72)
                            if text_info["font_bold"] is not None:
                                paragraph.font.bold = text_info["font_bold"]
                            if text_info["font_italic"] is not None:
                                paragraph.font.italic = text_info["font_italic"]
                            if text_info["font_color_rgb"] is not None:
                                rgb_int = text_info["font_color_rgb"]
                                r = (rgb_int >> 16) & 0xFF
                                g = (rgb_int >> 8) & 0xFF
                                b = rgb_int & 0xFF
                                paragraph.font.color.rgb = RGBColor(r, g, b)
                            # Alignment for table cells is also complex, skipping for now

    prs.save(output_ppt_path)

if __name__ == "__main__":
    # This is a simplified example. In a real application, you would get
    # translated_data from your LLM translation module.

    # Create a dummy PPT for testing
    prs = Presentation()
    slide_layout = prs.slide_layouts[1] # Title and Content layout

    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "Hello World!"
    content.text = "This is a test paragraph.\nAnother line of text.\n\nAnd a third one."

    # Add a table
    rows = 2
    cols = 2
    left = top = width = height = Inches(1.0)
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    cell = table.cell(0, 0)
    cell.text = "Header 1"
    cell = table.cell(0, 1)
    cell.text = "Header 2"
    cell = table.cell(1, 0)
    cell.text = "Data A"
    cell = table.cell(1, 1)
    cell.text = "Data B"

    # Add another slide with just a text box
    slide_layout_blank = prs.slide_layouts[6] # Blank layout
    slide2 = prs.slides.add_slide(slide_layout_blank)
    left = top = width = height = Inches(1.0)
    txBox = slide2.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "This is a standalone text box."

    dummy_ppt_path = "dummy_original.pptx"
    prs.save(dummy_ppt_path)
    print(f"Created dummy original PPTX at {dummy_ppt_path}")

    # Simulate extracted and translated data
    # In a real scenario, this would come from ppt_text_extractor and llm_translator
    simulated_translated_data = [
        {
            "slide_index": 0,
            "texts": [
                {
                    "shape_index": 0,
                    "paragraph_index": 0,
                    "original_text": "Hello World!",
                    "translated_text": "你好世界！",
                    "shape_type": "PLACEHOLDER",
                    "left": 457200,
                    "top": 274638,
                    "width": 8229600,
                    "height": 1143000,
                    "font_name": "Arial",
                    "font_size": 2438400,
                    "font_bold": True,
                    "font_italic": False,
                    "font_color_rgb": 0,
                    "alignment": "CENTER"
                },
                {
                    "shape_index": 1,
                    "paragraph_index": 0,
                    "original_text": "This is a test paragraph.",
                    "translated_text": "这是一个测试段落。",
                    "shape_type": "PLACEHOLDER",
                    "left": 457200,
                    "top": 1600200,
                    "width": 8229600,
                    "height": 4525963,
                    "font_name": "Calibri",
                    "font_size": 1828800,
                    "font_bold": False,
                    "font_italic": False,
                    "font_color_rgb": 0,
                    "alignment": "LEFT"
                },
                {
                    "shape_index": 1,
                    "paragraph_index": 1,
                    "original_text": "Another line of text.",
                    "translated_text": "另一行文字。",
                    "shape_type": "PLACEHOLDER",
                    "left": 457200,
                    "top": 1600200,
                    "width": 8229600,
                    "height": 4525963,
                    "font_name": "Calibri",
                    "font_size": 1828800,
                    "font_bold": False,
                    "font_italic": False,
                    "font_color_rgb": 0,
                    "alignment": "LEFT"
                },
                {
                    "shape_index": 1,
                    "paragraph_index": 3,
                    "original_text": "And a third one.",
                    "translated_text": "还有第三个。",
                    "shape_type": "PLACEHOLDER",
                    "left": 457200,
                    "top": 1600200,
                    "width": 8229600,
                    "height": 4525963,
                    "font_name": "Calibri",
                    "font_size": 1828800,
                    "font_bold": False,
                    "font_italic": False,
                    "font_color_rgb": 0,
                    "alignment": "LEFT"
                },
                {
                    "shape_index": 2,
                    "row_index": 0,
                    "col_index": 0,
                    "original_text": "Header 1",
                    "translated_text": "标题 1",
                    "shape_type": "TABLE_CELL",
                    "font_name": "Calibri",
                    "font_size": 1828800,
                    "font_bold": True,
                    "font_italic": False,
                    "font_color_rgb": 0,
                    "alignment": "CENTER"
                },
                {
                    "shape_index": 2,
                    "row_index": 0,
                    "col_index": 1,
                    "original_text": "Header 2",
                    "translated_text": "标题 2",
                    "shape_type": "TABLE_CELL",
                    "font_name": "Calibri",
                    "font_size": 1828800,
                    "font_bold": True,
                    "font_italic": False,
                    "font_color_rgb": 0,
                    "alignment": "CENTER"
                },
                {
                    "shape_index": 2,
                    "row_index": 1,
                    "col_index": 0,
                    "original_text": "Data A",
                    "translated_text": "数据 A",
                    "shape_type": "TABLE_CELL",
                    "font_name": "Calibri",
                    "font_size": 1828800,
                    "font_bold": False,
                    "font_italic": False,
                    "font_color_rgb": 0,
                    "alignment": "LEFT"
                },
                {
                    "shape_index": 2,
                    "row_index": 1,
                    "col_index": 1,
                    "original_text": "Data B",
                    "translated_text": "数据 B",
                    "shape_type": "TABLE_CELL",
                    "font_name": "Calibri",
                    "font_size": 1828800,
                    "font_bold": False,
                    "font_italic": False,
                    "font_color_rgb": 0,
                    "alignment": "LEFT"
                }
            ]
        },
        {
            "slide_index": 1,
            "texts": [
                {
                    "shape_index": 0,
                    "paragraph_index": 0,
                    "original_text": "This is a standalone text box.",
                    "translated_text": "这是一个独立的文本框。",
                    "shape_type": "TEXT_BOX",
                    "left": 914400,
                    "top": 914400,
                    "width": 914400,
                    "height": 914400,
                    "font_name": "Calibri",
                    "font_size": 1828800,
                    "font_bold": False,
                    "font_italic": False,
                    "font_color_rgb": 0,
                    "alignment": "LEFT"
                }
            ]
        }
    ]

    output_ppt_path = "dummy_translated.pptx"
    reconstruct_ppt(dummy_ppt_path, simulated_translated_data, output_ppt_path)
    print(f"Created translated PPTX at {output_ppt_path}")

    # Clean up dummy files
    import os
    os.remove(dummy_ppt_path)
    print(f"Removed dummy original PPTX at {dummy_ppt_path}")


